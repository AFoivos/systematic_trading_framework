from __future__ import annotations

from pathlib import Path
from typing import Any, Literal, Mapping


def write_market_making_presentation(
    run_dir: str | Path,
    diagnostics: Mapping[str, Any],
    *,
    output_path: str | Path | None = None,
    language: Literal["en", "el"] = "el",
) -> Path:
    """Write a PowerPoint diagnostics deck when python-pptx is available."""
    run_path = Path(run_dir)
    out = Path(output_path) if output_path is not None else run_path / "diagnostics" / "market_making_diagnostics.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        diagnostics.setdefault("warnings", []).append(f"PowerPoint generation skipped: python-pptx unavailable ({exc}).")
        return out

    labels = _labels(language)
    prs = Presentation()

    def add_text_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for idx, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0

    run = diagnostics.get("run", {})
    gaps = diagnostics.get("gaps", {})
    artifacts = diagnostics.get("artifacts", {})
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = labels["title"]
    title.placeholders[1].text = f"{run.get('data_source', 'n/a')} | {run.get('fill_model', 'n/a')} | {run.get('run_start_timestamp', 'n/a')}"
    add_text_slide(
        labels["summary"],
        [
            f"PnL: {run.get('total_pnl')}",
            f"Fills: {run.get('number_of_fills')}",
            f"Quotes: {run.get('number_of_quotes')}",
            f"Fill ratio: {run.get('fill_ratio')}",
            f"Fills/quote attempt: {run.get('fills_per_quote_attempt')}",
            f"Fills/placed quote: {run.get('fills_per_placed_quote')}",
            f"Drawdown: {run.get('max_drawdown')}",
            f"Kill switches: {run.get('kill_switch_events')}",
            f"Lineage available: {not gaps.get('lineage_missing', True)}",
            f"Markout available: {not gaps.get('markout_missing', True)}",
            f"Adverse-selection filter active: {run.get('adverse_selection_filter_active')}",
        ],
    )
    add_text_slide(labels["funnel"], [f"Input events: {run.get('input_events')}", f"Quoted events: {run.get('quoted_events')}", f"Fills: {run.get('number_of_fills')}"])
    _add_image_slide(prs, labels["pnl"], artifacts.get("pnl_curve.png"), artifacts.get("drawdown.png"), Inches)
    _add_image_slide(prs, labels["inventory"], artifacts.get("inventory_timeseries.png"), None, Inches)
    _add_image_slide(prs, labels["quotes"], artifacts.get("quoted_spread_distribution.png"), artifacts.get("book_spread_distribution.png"), Inches)
    _add_image_slide(prs, labels["fills"], artifacts.get("fill_side_counts.png"), None, Inches)
    _add_image_slide(prs, labels["markout"], artifacts.get("fill_markout_distribution.png"), None, Inches)
    add_text_slide(labels["quality"], [f"Crossed books: {diagnostics.get('market_quality', {}).get('crossed_book_count')}", f"Missing top of book: {diagnostics.get('market_quality', {}).get('missing_top_of_book_count')}"])
    add_text_slide(labels["gaps"], [f"{key}: {value}" for key, value in gaps.items()])
    add_text_slide(labels["recommendations"], ["Add queue-position, partial-fill, and latency models before live trading.", "Use markout and risk diagnostics to tune quote placement."])
    prs.save(out)
    return out


def _add_image_slide(prs: Any, title: str, image1: str | None, image2: str | None, Inches: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image1 and Path(image1).exists():
        slide.shapes.add_picture(image1, Inches(0.7), Inches(1.4), width=Inches(4.2))
    if image2 and Path(image2).exists():
        slide.shapes.add_picture(image2, Inches(5.0), Inches(1.4), width=Inches(4.2))


def _labels(language: Literal["en", "el"]) -> dict[str, str]:
    if language == "en":
        return {
            "title": "Market Making Diagnostics",
            "summary": "Executive summary",
            "funnel": "Event funnel",
            "pnl": "PnL / drawdown",
            "inventory": "Inventory behavior",
            "quotes": "Quote behavior",
            "fills": "Fill behavior",
            "markout": "Markout / adverse selection",
            "quality": "Market quality",
            "gaps": "Diagnostics gaps",
            "recommendations": "Recommendations",
        }
    return {
        "title": "Διαγνωστικά Market Making",
        "summary": "Σύνοψη",
        "funnel": "Ροή γεγονότων",
        "pnl": "PnL / drawdown",
        "inventory": "Συμπεριφορά inventory",
        "quotes": "Συμπεριφορά quotes",
        "fills": "Συμπεριφορά fills",
        "markout": "Markout / adverse selection",
        "quality": "Ποιότητα order book",
        "gaps": "Κενά διαγνωστικών",
        "recommendations": "Προτεινόμενα επόμενα βήματα",
    }


__all__ = ["write_market_making_presentation"]
