from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Mapping, Sequence

import numpy as np
import pandas as pd


DIAGNOSTIC_KEYS = [
    "run",
    "quote",
    "fill",
    "pnl",
    "inventory",
    "market_quality",
    "risk",
    "markout",
    "gaps",
    "warnings",
    "artifacts",
]


def build_market_making_diagnostics(
    run_dir: str | Path,
    *,
    orderbook_events_path: str | Path | None = None,
    markout_horizons: Sequence[int] = (1, 5, 10, 30),
    max_inventory: float | None = None,
) -> dict[str, Any]:
    """Build tolerant diagnostics for a local market-making run directory."""
    run_path = Path(run_dir)
    if not run_path.exists() or not run_path.is_dir():
        raise FileNotFoundError(f"Market-making run_dir does not exist: {run_path}")

    warnings: list[str] = []
    artifacts: dict[str, str] = {}
    summary = _read_json(run_path / "summary.json", warnings=warnings)
    orders = _read_csv(run_path / "orders.csv", warnings=warnings, optional=True)
    trades = _read_csv(run_path / "trades.csv", warnings=warnings, optional=True)
    pnl = _read_csv(run_path / "pnl_timeseries.csv", warnings=warnings, optional=True)
    inventory = _read_csv(run_path / "inventory_timeseries.csv", warnings=warnings, optional=True)
    quote_events = _read_csv(run_path / "quote_events.csv", warnings=warnings, optional=True)

    if orderbook_events_path is not None:
        orderbook_path = Path(orderbook_events_path)
    else:
        orderbook_path = run_path / "orderbook_events.csv"
    orderbook = _read_csv(orderbook_path, warnings=warnings, optional=True)

    for name in [
        "summary.json",
        "orders.csv",
        "trades.csv",
        "pnl_timeseries.csv",
        "inventory_timeseries.csv",
        "quote_events.csv",
    ]:
        path = run_path / name
        if path.exists():
            artifacts[name] = str(path)
    if orderbook_path.exists():
        artifacts["orderbook_events.csv"] = str(orderbook_path)

    quote_diag, quote_rows = _quote_diagnostics(quote_events, warnings)
    fill_diag, fill_rows = _fill_diagnostics(trades, orders, quote_events, warnings)
    pnl_diag, pnl_rows = _pnl_diagnostics(summary, pnl, trades, orders, inventory, warnings)
    inventory_diag, inventory_rows = _inventory_diagnostics(inventory, max_inventory, warnings)
    market_diag, market_rows = _market_quality_diagnostics(orderbook, warnings)
    risk_diag, risk_rows = _risk_diagnostics(quote_events, summary, warnings)
    markout_diag, markout_rows = _markout_diagnostics(trades, orderbook, markout_horizons, warnings)
    gaps = _gap_diagnostics(
        quote_events=quote_events,
        orders=orders,
        orderbook=orderbook,
        fills=trades,
        markout_rows=markout_rows,
    )
    run_diag = _run_diagnostics(
        summary=summary,
        quote_diag=quote_diag,
        fill_diag=fill_diag,
        pnl_diag=pnl_diag,
        inventory_diag=inventory_diag,
        pnl=pnl,
        inventory=inventory,
        orderbook=orderbook,
        gaps=gaps,
        warnings=warnings,
    )

    diagnostics: dict[str, Any] = {
        "run": run_diag,
        "quote": quote_diag,
        "fill": fill_diag,
        "pnl": pnl_diag,
        "inventory": inventory_diag,
        "market_quality": market_diag,
        "risk": risk_diag,
        "markout": markout_diag,
        "gaps": gaps,
        "warnings": warnings,
        "artifacts": artifacts,
        "_tables": {
            "quote_diagnostics": quote_rows,
            "fill_diagnostics": fill_rows,
            "pnl_attribution": pnl_rows,
            "inventory_diagnostics": inventory_rows,
            "market_quality": market_rows,
            "risk_diagnostics": risk_rows,
            "markout_diagnostics": markout_rows,
        },
    }
    return diagnostics


def write_market_making_diagnostics(
    run_dir: str | Path,
    *,
    orderbook_events_path: str | Path | None = None,
    diagnostics_dir: str | Path | None = None,
    markout_horizons: Sequence[int] = (1, 5, 10, 30),
    max_inventory: float | None = None,
    make_plots: bool = True,
) -> dict[str, Any]:
    """Build diagnostics and write local JSON/CSV/plot artifacts."""
    run_path = Path(run_dir)
    out_dir = Path(diagnostics_dir) if diagnostics_dir is not None else run_path / "diagnostics"
    out_dir.mkdir(parents=True, exist_ok=True)
    diagnostics = build_market_making_diagnostics(
        run_path,
        orderbook_events_path=orderbook_events_path,
        markout_horizons=markout_horizons,
        max_inventory=max_inventory,
    )
    tables = diagnostics.pop("_tables")

    for name, frame in tables.items():
        if isinstance(frame, pd.DataFrame) and not frame.empty:
            path = out_dir / f"{name}.csv"
            frame.to_csv(path, index=False)
            diagnostics["artifacts"][f"{name}.csv"] = str(path)

    if make_plots:
        _write_plots(out_dir, tables=tables, diagnostics=diagnostics)

    summary_path = out_dir / "summary.json"
    gaps_path = out_dir / "gaps.json"
    diagnostics["artifacts"]["diagnostics_summary"] = str(summary_path)
    diagnostics["artifacts"]["gaps"] = str(gaps_path)
    summary_path.write_text(json.dumps(_jsonable(diagnostics), indent=2), encoding="utf-8")
    gaps_path.write_text(json.dumps(_jsonable(diagnostics["gaps"]), indent=2), encoding="utf-8")
    return diagnostics


def discover_market_making_runs(reports_root: str | Path = "reports") -> list[Path]:
    """Find report directories containing market-making run artifacts."""
    root = Path(reports_root)
    if not root.exists():
        return []
    markers = {"summary.json", "orders.csv", "trades.csv", "pnl_timeseries.csv", "inventory_timeseries.csv"}
    runs: list[Path] = []
    for path in [root, *[p for p in root.rglob("*") if p.is_dir()]]:
        if path.name == "diagnostics" or "comparison" in path.name:
            continue
        if any(part == "diagnostics" for part in path.parts):
            continue
        names = {child.name for child in path.iterdir() if child.is_file()}
        if markers & names:
            runs.append(path)
    return sorted(set(runs), key=lambda p: p.stat().st_mtime, reverse=True)


def write_market_making_comparison(runs: Sequence[Path], output_dir: str | Path) -> dict[str, Path]:
    """Write a compact multi-run comparison summary and markdown report."""
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    rows: list[dict[str, Any]] = []
    for run in runs:
        diagnostics = build_market_making_diagnostics(run)
        rows.append(
            {
                "run_dir": str(run),
                "total_pnl": diagnostics["run"].get("total_pnl"),
                "fills": diagnostics["run"].get("number_of_fills"),
                "quotes": diagnostics["run"].get("number_of_quotes"),
                "fill_ratio": diagnostics["run"].get("fill_ratio"),
                "max_drawdown": diagnostics["run"].get("max_drawdown"),
                "fees": diagnostics["run"].get("fees"),
                "avg_spread": diagnostics["run"].get("average_spread_quoted"),
                "avg_abs_inventory": diagnostics["inventory"].get("avg_abs_inventory"),
                "kill_switches": len(diagnostics["run"].get("kill_switch_events") or []),
                "warnings_count": len(diagnostics["warnings"]),
                "diagnostics_gaps": sum(1 for v in diagnostics["gaps"].values() if bool(v) is True),
            }
        )
    frame = pd.DataFrame(rows)
    summary_path = out / "summary.csv"
    report_path = out / "report.md"
    pptx_path = out / "comparison.pptx"
    frame.to_csv(summary_path, index=False)
    report_path.write_text(_comparison_markdown(frame), encoding="utf-8")
    _write_comparison_pptx(frame, pptx_path)
    return {"summary": summary_path, "report": report_path, "pptx": pptx_path}


def _run_diagnostics(
    *,
    summary: Mapping[str, Any],
    quote_diag: Mapping[str, Any],
    fill_diag: Mapping[str, Any],
    pnl_diag: Mapping[str, Any],
    inventory_diag: Mapping[str, Any],
    pnl: pd.DataFrame,
    inventory: pd.DataFrame,
    orderbook: pd.DataFrame,
    gaps: Mapping[str, Any],
    warnings: list[str],
) -> dict[str, Any]:
    run_start, run_end = _run_bounds([pnl, inventory, orderbook])
    duration = (run_end - run_start).total_seconds() if run_start is not None and run_end is not None else None
    input_events = _coalesce(summary.get("input_events"), len(orderbook) if not orderbook.empty else None)
    events_per_second = _safe_div(input_events, duration) if duration and input_events is not None else None
    if duration is None:
        warnings.append("run_start_timestamp/run_end_timestamp unavailable.")
    warnings.append("fill_ratio is kept for backward compatibility and is currently an alias of fills_per_order.")
    return {
        "total_pnl": _coalesce(summary.get("total_pnl"), pnl_diag.get("net_pnl")),
        "realized_pnl": _coalesce(summary.get("realized_pnl"), pnl_diag.get("realized_pnl_final")),
        "unrealized_pnl": _coalesce(summary.get("unrealized_pnl"), pnl_diag.get("unrealized_pnl_final")),
        "fees": _coalesce(summary.get("fees"), fill_diag.get("total_fees")),
        "fee_drag_ratio": pnl_diag.get("fee_drag_ratio"),
        "number_of_quotes": _coalesce(summary.get("number_of_quotes"), quote_diag.get("placed_quote_count")),
        "number_of_fills": _coalesce(summary.get("number_of_fills"), fill_diag.get("fill_count")),
        "number_of_cancels": summary.get("number_of_cancels"),
        "fill_ratio": _coalesce(summary.get("fill_ratio"), fill_diag.get("fills_per_order")),
        "fills_per_quote_attempt": _coalesce(summary.get("fills_per_quote_attempt"), _safe_div(fill_diag.get("fill_count"), quote_diag.get("quote_count"))),
        "fills_per_placed_quote": _coalesce(summary.get("fills_per_placed_quote"), _safe_div(fill_diag.get("fill_count"), quote_diag.get("placed_quote_count"))),
        "fills_per_order": _coalesce(summary.get("fills_per_order"), fill_diag.get("fills_per_order")),
        "fills_per_input_event": _coalesce(summary.get("fills_per_input_event"), _safe_div(fill_diag.get("fill_count"), input_events)),
        "cancel_to_quote_ratio": _safe_div(summary.get("number_of_cancels"), summary.get("number_of_quotes")),
        "quoted_events": _coalesce(summary.get("quoted_events"), quote_diag.get("placed_quote_count")),
        "skipped_events": summary.get("skipped_events"),
        "input_events": input_events,
        "quoted_event_rate": _safe_div(summary.get("quoted_events"), input_events),
        "skipped_event_rate": _safe_div(summary.get("skipped_events"), input_events),
        "average_spread_quoted": _coalesce(summary.get("average_spread_quoted"), quote_diag.get("avg_quoted_spread_bps")),
        "average_inventory": _coalesce(summary.get("average_inventory"), inventory_diag.get("avg_inventory")),
        "average_abs_inventory": inventory_diag.get("avg_abs_inventory"),
        "max_inventory": _coalesce(summary.get("max_inventory"), inventory_diag.get("max_abs_inventory")),
        "max_drawdown": _coalesce(summary.get("max_drawdown"), pnl_diag.get("max_drawdown")),
        "kill_switch_events": summary.get("kill_switch_events", []),
        "runtime_errors": summary.get("runtime_errors"),
        "reconnects": summary.get("reconnects"),
        "fill_model": summary.get("fill_model"),
        "data_source": summary.get("data_source"),
        "quote_event_logging_available": not gaps.get("quote_events_missing", True),
        "markout_available": not gaps.get("markout_missing", True),
        "lineage_available": not gaps.get("lineage_missing", True),
        "pnl_attribution_approximate": gaps.get("pnl_attribution_approximate", True),
        "adverse_selection_filter_active": bool(summary.get("adverse_selection_filter_active", False)),
        "run_start_timestamp": run_start.isoformat() if run_start is not None else None,
        "run_end_timestamp": run_end.isoformat() if run_end is not None else None,
        "run_duration_seconds": duration,
        "events_per_second": events_per_second,
    }


def _quote_diagnostics(frame: pd.DataFrame, warnings: list[str]) -> tuple[dict[str, Any], pd.DataFrame]:
    if frame.empty:
        warnings.append("quote_events.csv missing; quote diagnostics are limited.")
        return {"quote_event_logging_missing": True}, pd.DataFrame()
    q = frame.copy()
    _coerce_numeric(q)
    for col in ["placed", "should_quote", "risk_allowed", "risk_cancel_all", "risk_kill_switch"]:
        if col in q:
            q[col] = _to_bool(q[col])
    q["quote_mid_price"] = (q["bid_price"] + q["ask_price"]) / 2.0
    q["quote_spread"] = q["ask_price"] - q["bid_price"]
    q["quote_spread_bps"] = np.where(q["quote_mid_price"] > 0, q["quote_spread"] / q["quote_mid_price"] * 10_000.0, np.nan)
    q["bid_distance_from_fair_bps"] = np.where(q["fair_price"] > 0, (q["fair_price"] - q["bid_price"]) / q["fair_price"] * 10_000.0, np.nan)
    q["ask_distance_from_fair_bps"] = np.where(q["fair_price"] > 0, (q["ask_price"] - q["fair_price"]) / q["fair_price"] * 10_000.0, np.nan)
    q["bid_distance_from_book_bps"] = np.where(q["book_best_bid"] > 0, (q["bid_price"] - q["book_best_bid"]) / q["book_best_bid"] * 10_000.0, np.nan)
    q["ask_distance_from_book_bps"] = np.where(q["book_best_ask"] > 0, (q["book_best_ask"] - q["ask_price"]) / q["book_best_ask"] * 10_000.0, np.nan)
    q["quote_inside_book"] = (q["bid_price"] > q["book_best_bid"]) & (q["ask_price"] < q["book_best_ask"])
    q["quote_crosses_book"] = (q["bid_price"] >= q["book_best_ask"]) | (q["ask_price"] <= q["book_best_bid"])
    q["quote_at_top_of_book"] = (q["bid_price"] == q["book_best_bid"]) & (q["ask_price"] == q["book_best_ask"])
    q["abs_inventory_ratio"] = q.get("inventory_ratio", pd.Series(dtype=float)).abs()
    quote_count = len(q)
    placed = int(q.get("placed", pd.Series(False, index=q.index)).sum())
    rejected = quote_count - placed
    return {
        "quote_event_logging_missing": False,
        "quote_count": quote_count,
        "placed_quote_count": placed,
        "rejected_quote_count": rejected,
        "placed_quote_rate": _safe_div(placed, quote_count),
        "rejected_quote_rate": _safe_div(rejected, quote_count),
        "reject_count_by_reason": _counts(q.loc[~q.get("placed", False), "quote_reason"] if "quote_reason" in q else pd.Series(dtype=str)),
        "risk_reject_count_by_reason": _counts(q.loc[~q.get("risk_allowed", False), "risk_reason"] if "risk_reason" in q else pd.Series(dtype=str)),
        "avg_quoted_spread_bps": _mean(q["quote_spread_bps"]),
        "median_quoted_spread_bps": _median(q["quote_spread_bps"]),
        "p95_quoted_spread_bps": _quantile(q["quote_spread_bps"], 0.95),
        "min_quoted_spread_bps": _min(q["quote_spread_bps"]),
        "max_quoted_spread_bps": _max(q["quote_spread_bps"]),
        "avg_book_spread_bps": _mean(q.get("book_spread_bps")),
        "median_book_spread_bps": _median(q.get("book_spread_bps")),
        "p95_book_spread_bps": _quantile(q.get("book_spread_bps"), 0.95),
        "quote_vs_book_spread_ratio": _safe_div(_mean(q["quote_spread_bps"]), _mean(q.get("book_spread_bps"))),
        "avg_inventory_ratio": _mean(q.get("inventory_ratio")),
        "p95_abs_inventory_ratio": _quantile(q["abs_inventory_ratio"], 0.95),
        "avg_book_imbalance_1": _mean(q.get("book_imbalance_1")),
        "avg_book_imbalance_5": _mean(q.get("book_imbalance_5")),
        "quote_cross_count": int(q["quote_crosses_book"].sum()),
        "quote_inside_book_count": int(q["quote_inside_book"].sum()),
        "quote_at_top_of_book_count": int(q["quote_at_top_of_book"].sum()),
    }, q


def _fill_diagnostics(
    trades: pd.DataFrame,
    orders: pd.DataFrame,
    quote_events: pd.DataFrame,
    warnings: list[str],
) -> tuple[dict[str, Any], pd.DataFrame]:
    if trades.empty:
        warnings.append("trades.csv missing or empty; fill diagnostics are limited.")
        return {"fill_count": 0}, pd.DataFrame()
    fills = trades.copy()
    _coerce_numeric(fills)
    fills["timestamp"] = _to_datetime(fills.get("timestamp"))
    fills["fill_notional"] = fills["price"] * fills["quantity"]
    if not orders.empty and "order_id" in orders and "order_id" in fills:
        order_cols = ["order_id", "timestamp"]
        if "parent_quote_event_id" in orders:
            order_cols.append("parent_quote_event_id")
        order_times = orders[order_cols].copy()
        order_times["order_timestamp"] = _to_datetime(order_times["timestamp"])
        merge_cols = ["order_id", "order_timestamp"]
        if "parent_quote_event_id" in order_times and "parent_quote_event_id" not in fills:
            merge_cols.append("parent_quote_event_id")
        fills = fills.merge(order_times[merge_cols], on="order_id", how="left")
        fills["time_to_fill_seconds"] = (fills["timestamp"] - fills["order_timestamp"]).dt.total_seconds()
    else:
        fills["time_to_fill_seconds"] = np.nan
        warnings.append("orders.csv cannot be joined to trades.csv; time-to-fill metrics unavailable.")
    if not quote_events.empty and "quote_event_id" in quote_events and "parent_quote_event_id" in fills:
        quote_cols = [
            col
            for col in ["quote_event_id", "quote_reason", "risk_reason", "spread_bps"]
            if col in quote_events
        ]
        fills = fills.merge(
            quote_events[quote_cols].rename(columns={"quote_event_id": "parent_quote_event_id", "spread_bps": "quote_spread_bps"}),
            on="parent_quote_event_id",
            how="left",
        )
        if "quote_spread_bps" in fills:
            fills["quote_spread_bucket"] = pd.cut(
                fills["quote_spread_bps"],
                bins=[-np.inf, 1, 5, 10, 25, np.inf],
                labels=["<=1", "1-5", "5-10", "10-25", ">25"],
            ).astype(str)
    fill_count = len(fills)
    buy_count = int((fills["side"] == "buy").sum()) if "side" in fills else 0
    sell_count = int((fills["side"] == "sell").sum()) if "side" in fills else 0
    filled_quote_count = int(fills["parent_quote_event_id"].dropna().nunique()) if "parent_quote_event_id" in fills else None
    return {
        "fill_count": fill_count,
        "buy_fill_count": buy_count,
        "sell_fill_count": sell_count,
        "buy_fill_rate": _safe_div(buy_count, fill_count),
        "sell_fill_rate": _safe_div(sell_count, fill_count),
        "avg_fill_price": _mean(fills.get("price")),
        "avg_fill_notional": _mean(fills["fill_notional"]),
        "median_fill_notional": _median(fills["fill_notional"]),
        "total_fill_notional": _sum(fills["fill_notional"]),
        "total_fees": _sum(fills.get("fee")),
        "avg_fee_per_fill": _mean(fills.get("fee")),
        "fills_per_quote": None,
        "fills_per_order": _safe_div(fill_count, len(orders) if not orders.empty else None),
        "filled_quote_count": filled_quote_count,
        "filled_quote_rate": _safe_div(filled_quote_count, orders["parent_quote_event_id"].nunique() if "parent_quote_event_id" in orders else None),
        "fills_by_quote_reason": _counts(fills.get("quote_reason")),
        "fills_by_risk_reason": _counts(fills.get("risk_reason")),
        "fills_by_spread_bucket": _counts(fills.get("quote_spread_bucket")),
        "time_to_fill_seconds_mean": _mean(fills["time_to_fill_seconds"]),
        "time_to_fill_seconds_median": _median(fills["time_to_fill_seconds"]),
        "time_to_fill_seconds_p95": _quantile(fills["time_to_fill_seconds"], 0.95),
        "time_to_fill_by_quote_event_id_mean": _mean(fills.groupby("parent_quote_event_id")["time_to_fill_seconds"].min() if "parent_quote_event_id" in fills else None),
        "events_to_fill_mean": None,
        "events_to_fill_median": None,
    }, fills[
        [
            c
            for c in [
                "order_id",
                "parent_quote_event_id",
                "symbol",
                "side",
                "price",
                "quantity",
                "fee",
                "timestamp",
                "fill_notional",
                "time_to_fill_seconds",
                "quote_reason",
                "risk_reason",
                "quote_spread_bps",
                "quote_spread_bucket",
            ]
            if c in fills
        ]
    ]


def _markout_diagnostics(
    trades: pd.DataFrame,
    orderbook: pd.DataFrame,
    horizons: Sequence[int],
    warnings: list[str],
) -> tuple[dict[str, Any], pd.DataFrame]:
    if trades.empty or orderbook.empty:
        warnings.append("Markout unavailable: trades.csv or orderbook_events.csv missing.")
        return {"markout_unavailable": True}, pd.DataFrame()
    fills = trades.copy()
    book = orderbook.copy()
    _coerce_numeric(fills)
    _coerce_numeric(book)
    fills["timestamp"] = _to_datetime(fills.get("timestamp"))
    book["timestamp"] = _to_datetime(book.get("timestamp"))
    book = book.sort_values("timestamp").reset_index(drop=True)
    if "mid_price" not in book:
        warnings.append("Markout unavailable: orderbook mid_price missing.")
        return {"markout_unavailable": True}, pd.DataFrame()
    rows: list[dict[str, Any]] = []
    for _, fill in fills.sort_values("timestamp").iterrows():
        future_start = int(book["timestamp"].searchsorted(fill["timestamp"], side="right"))
        row = fill.to_dict()
        for horizon in horizons:
            idx = future_start + int(horizon) - 1
            value = np.nan
            if 0 <= idx < len(book):
                future_mid = book.loc[idx, "mid_price"]
                if fill["side"] == "buy":
                    value = (future_mid - fill["price"]) / fill["price"] * 10_000.0
                elif fill["side"] == "sell":
                    value = (fill["price"] - future_mid) / fill["price"] * 10_000.0
            row[f"markout_bps_h{horizon}"] = value
        rows.append(row)
    markouts = pd.DataFrame(rows)
    result: dict[str, Any] = {"markout_unavailable": False}
    for horizon in horizons:
        col = f"markout_bps_h{horizon}"
        result[f"avg_markout_bps_h{horizon}"] = _mean(markouts.get(col))
        result[f"median_markout_bps_h{horizon}"] = _median(markouts.get(col))
        negative_rate = _mean((markouts[col] < 0).astype(float)) if col in markouts and markouts[col].notna().any() else None
        result[f"negative_markout_rate_h{horizon}"] = negative_rate
        result[f"adverse_selection_rate_h{horizon}"] = negative_rate
    return result, markouts


def _pnl_diagnostics(
    summary: Mapping[str, Any],
    pnl: pd.DataFrame,
    trades: pd.DataFrame,
    orders: pd.DataFrame,
    inventory: pd.DataFrame,
    warnings: list[str],
) -> tuple[dict[str, Any], pd.DataFrame]:
    if pnl.empty:
        warnings.append("pnl_timeseries.csv missing; PnL diagnostics use summary only.")
        net = summary.get("total_pnl")
        return {"net_pnl": net, "pnl_attribution_approximate": True}, pd.DataFrame()
    p = pnl.copy()
    _coerce_numeric(p)
    p["timestamp"] = _to_datetime(p.get("timestamp"))
    total = p.get("total_pnl", pd.Series(dtype=float))
    fees = _coalesce(summary.get("fees"), _last(p.get("fees")))
    net = _last(total)
    drawdowns = total - total.cummax()
    max_dd = abs(float(drawdowns.min())) if not drawdowns.empty else None
    max_dd_ts = p.loc[drawdowns.idxmin(), "timestamp"].isoformat() if not drawdowns.empty and drawdowns.notna().any() else None
    notional = None
    if not trades.empty:
        t = trades.copy()
        _coerce_numeric(t)
        notional = _sum(t.get("price") * t.get("quantity"))
    lineage_available = (
        not trades.empty
        and not orders.empty
        and "parent_quote_event_id" in trades
        and "parent_quote_event_id" in orders
    )
    if not lineage_available:
        warnings.append("PnL attribution is approximate unless full quote/order/fill linkage is available.")
    result = {
        "net_pnl": net,
        "gross_pnl_estimate": net + fees if net is not None and fees is not None else None,
        "realized_pnl_final": _last(p.get("realized_pnl")),
        "unrealized_pnl_final": _last(p.get("unrealized_pnl")),
        "fees": fees,
        "fee_drag_ratio": _safe_div(fees, abs(net) if net is not None else None),
        "spread_capture_estimate": None,
        "inventory_pnl_estimate": _last(p.get("unrealized_pnl")),
        "realized_pnl_change": _delta(p.get("realized_pnl")),
        "unrealized_pnl_change": _delta(p.get("unrealized_pnl")),
        "pnl_per_fill": _safe_div(net, len(trades) if not trades.empty else None),
        "pnl_per_quote": _safe_div(net, summary.get("number_of_quotes")),
        "pnl_per_order": _safe_div(net, len(orders) if not orders.empty else None),
        "pnl_per_notional": _safe_div(net, notional),
        "max_drawdown": max_dd,
        "max_drawdown_timestamp": max_dd_ts,
        "positive_pnl_time_rate": _mean((total > 0).astype(float)),
        "negative_pnl_time_rate": _mean((total < 0).astype(float)),
        "pnl_attribution_approximate": not lineage_available,
    }
    return result, pd.DataFrame([result])


def _inventory_diagnostics(frame: pd.DataFrame, max_inventory: float | None, warnings: list[str]) -> tuple[dict[str, Any], pd.DataFrame]:
    if frame.empty or "inventory" not in frame:
        warnings.append("inventory_timeseries.csv missing; inventory diagnostics unavailable.")
        return {}, pd.DataFrame()
    inv = frame.copy()
    _coerce_numeric(inv)
    inv["timestamp"] = _to_datetime(inv.get("timestamp"))
    series = inv["inventory"]
    abs_inv = series.abs()
    signs = np.sign(series.fillna(0.0))
    result = {
        "avg_inventory": _mean(series),
        "avg_abs_inventory": _mean(abs_inv),
        "median_abs_inventory": _median(abs_inv),
        "max_abs_inventory": _max(abs_inv),
        "p95_abs_inventory": _quantile(abs_inv, 0.95),
        "inventory_std": _std(series),
        "inventory_sign_flip_count": int(((signs.shift(1) * signs) < 0).sum()),
        "inventory_zero_cross_count": int(((series.shift(1).fillna(0) * series.fillna(0)) < 0).sum()),
        "time_weighted_abs_inventory": _mean(abs_inv),
        "pct_time_inventory_positive": _mean((series > 0).astype(float)),
        "pct_time_inventory_negative": _mean((series < 0).astype(float)),
        "pct_time_inventory_flat": _mean((series == 0).astype(float)),
    }
    if max_inventory is None or max_inventory <= 0:
        warnings.append("max_inventory unavailable; inventory limit utilization metrics skipped.")
        result.update(
            {
                "inventory_limit_utilization": None,
                "pct_time_above_25pct_limit": None,
                "pct_time_above_50pct_limit": None,
                "pct_time_above_75pct_limit": None,
                "pct_time_above_90pct_limit": None,
            }
        )
    else:
        utilization = abs_inv / float(max_inventory)
        result.update(
            {
                "inventory_limit_utilization": _max(utilization),
                "pct_time_above_25pct_limit": _mean((utilization > 0.25).astype(float)),
                "pct_time_above_50pct_limit": _mean((utilization > 0.50).astype(float)),
                "pct_time_above_75pct_limit": _mean((utilization > 0.75).astype(float)),
                "pct_time_above_90pct_limit": _mean((utilization > 0.90).astype(float)),
            }
        )
    return result, inv


def _market_quality_diagnostics(frame: pd.DataFrame, warnings: list[str]) -> tuple[dict[str, Any], pd.DataFrame]:
    if frame.empty:
        warnings.append("orderbook_events.csv missing; market-quality diagnostics unavailable.")
        return {"orderbook_events_missing": True}, pd.DataFrame()
    q = frame.copy()
    _coerce_numeric(q)
    q["timestamp"] = _to_datetime(q.get("timestamp"))
    q = q.sort_values("timestamp").reset_index(drop=True)
    q["timestamp_gap_ms"] = q["timestamp"].diff().dt.total_seconds() * 1000.0
    q["is_crossed_book"] = q["best_bid"] >= q["best_ask"]
    q["is_missing_top_of_book"] = q["best_bid"].isna() | q["best_ask"].isna()
    p95_spread = _quantile(q.get("spread_bps"), 0.95)
    q["is_wide_spread"] = q["spread_bps"] > p95_spread if p95_spread is not None else False
    seq_gap = 0
    if "sequence" in q and q["sequence"].notna().any():
        seq_gap = int((q["sequence"].dropna().diff() > 1).sum())
    return {
        "orderbook_events_missing": False,
        "event_count": len(q),
        "snapshot_count": int((q.get("event_type") == "snapshot").sum()) if "event_type" in q else None,
        "update_count": int((q.get("event_type") == "update").sum()) if "event_type" in q else None,
        "invalid_book_count": int(q["is_crossed_book"].sum() + q["is_missing_top_of_book"].sum()),
        "crossed_book_count": int(q["is_crossed_book"].sum()),
        "missing_top_of_book_count": int(q["is_missing_top_of_book"].sum()),
        "avg_spread_bps": _mean(q.get("spread_bps")),
        "median_spread_bps": _median(q.get("spread_bps")),
        "p95_spread_bps": p95_spread,
        "max_spread_bps": _max(q.get("spread_bps")),
        "avg_imbalance_1": _mean(q.get("imbalance_1")),
        "median_imbalance_1": _median(q.get("imbalance_1")),
        "p95_imbalance_1": _quantile(q.get("imbalance_1"), 0.95),
        "avg_imbalance_5": _mean(q.get("imbalance_5")),
        "median_imbalance_5": _median(q.get("imbalance_5")),
        "p95_imbalance_5": _quantile(q.get("imbalance_5"), 0.95),
        "avg_bid_depth_5": _mean(q.get("bid_depth_5")),
        "avg_ask_depth_5": _mean(q.get("ask_depth_5")),
        "median_bid_depth_5": _median(q.get("bid_depth_5")),
        "median_ask_depth_5": _median(q.get("ask_depth_5")),
        "timestamp_gap_avg_ms": _mean(q["timestamp_gap_ms"]),
        "timestamp_gap_median_ms": _median(q["timestamp_gap_ms"]),
        "timestamp_gap_p95_ms": _quantile(q["timestamp_gap_ms"], 0.95),
        "timestamp_gap_max_ms": _max(q["timestamp_gap_ms"]),
        "possible_sequence_gap_count": seq_gap,
        "duplicate_timestamp_count": int(q["timestamp"].duplicated().sum()),
    }, q[[c for c in ["timestamp", "symbol", "event_type", "best_bid", "best_ask", "mid_price", "spread", "spread_bps", "imbalance_1", "imbalance_5", "bid_depth_5", "ask_depth_5", "timestamp_gap_ms", "is_crossed_book", "is_missing_top_of_book", "is_wide_spread"] if c in q]]


def _risk_diagnostics(quote_events: pd.DataFrame, summary: Mapping[str, Any], warnings: list[str]) -> tuple[dict[str, Any], pd.DataFrame]:
    if quote_events.empty:
        warnings.append("quote_events.csv missing; using summary-only risk diagnostics.")
        return {
            "no_risk_reject_details": True,
            "kill_switch_events": summary.get("kill_switch_events", []),
        }, pd.DataFrame()
    q = quote_events.copy()
    for col in ["risk_allowed", "risk_cancel_all", "risk_kill_switch"]:
        if col in q:
            q[col] = _to_bool(q[col])
    reasons = q.get("risk_reason", pd.Series(dtype=str)).fillna("")
    result = {
        "no_risk_reject_details": False,
        "risk_allowed_count": int(q.get("risk_allowed", pd.Series(False, index=q.index)).sum()),
        "risk_reject_count": int((~q.get("risk_allowed", pd.Series(True, index=q.index))).sum()),
        "risk_cancel_all_count": int(q.get("risk_cancel_all", pd.Series(False, index=q.index)).sum()),
        "risk_kill_switch_count": int(q.get("risk_kill_switch", pd.Series(False, index=q.index)).sum()),
        "risk_reason_counts": _counts(reasons),
        "quote_reason_counts": _counts(q.get("quote_reason", pd.Series(dtype=str))),
        "kill_switch_events": summary.get("kill_switch_events", []),
        "stale_book_reject_count": _contains_count(reasons, "stale"),
        "extreme_spread_reject_count": _contains_count(reasons, "spread"),
        "max_inventory_reject_count": _contains_count(reasons, "max inventory"),
        "max_position_value_reject_count": _contains_count(reasons, "position value"),
        "max_daily_loss_reject_count": _contains_count(reasons, "daily loss"),
        "worst_case_inventory_reject_count": _contains_count(reasons, "worst-case inventory"),
        "worst_case_position_value_reject_count": _contains_count(reasons, "worst-case position"),
        "max_order_size_reject_count": _contains_count(reasons, "order size"),
        "max_open_orders_reject_count": _contains_count(reasons, "open orders"),
    }
    return result, pd.DataFrame([{"reason": k, "count": v} for k, v in result["risk_reason_counts"].items()])


def _gap_diagnostics(
    *,
    quote_events: pd.DataFrame,
    orders: pd.DataFrame,
    orderbook: pd.DataFrame,
    fills: pd.DataFrame,
    markout_rows: pd.DataFrame,
) -> dict[str, Any]:
    fill_count = 0 if fills.empty else len(fills)
    lineage_available = (
        not quote_events.empty
        and not orders.empty
        and (fills.empty or "parent_quote_event_id" in fills)
        and "quote_event_id" in quote_events
        and "parent_quote_event_id" in orders
    )
    return {
        "quote_events_missing": quote_events.empty,
        "markout_missing": markout_rows.empty,
        "orderbook_events_missing": orderbook.empty,
        "lineage_missing": not lineage_available,
        "fill_count_too_low_for_edge_evaluation": fill_count < 30,
        "no_risk_reject_details": quote_events.empty,
        "pnl_attribution_approximate": not lineage_available,
        "missing_queue_position_model": True,
        "missing_partial_fill_model": True,
        "missing_latency_model": True,
        "missing_adverse_selection_filter_diagnostics": True,
    }


def _write_plots(out_dir: Path, *, tables: Mapping[str, pd.DataFrame], diagnostics: dict[str, Any]) -> None:
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as exc:  # pragma: no cover - depends on optional local plotting stack
        diagnostics["warnings"].append(f"Plotting unavailable: {exc}")
        return

    def save_plot(name: str, fn: Any) -> None:
        try:
            fig = plt.figure(figsize=(8, 4.5))
            fn(plt)
            fig.tight_layout()
            path = out_dir / name
            fig.savefig(path)
            diagnostics["artifacts"][name] = str(path)
            plt.close(fig)
        except Exception as exc:  # pragma: no cover - defensive plotting isolation
            diagnostics["warnings"].append(f"Failed to create plot {name}: {exc}")
            plt.close("all")

    pnl = tables.get("pnl_attribution", pd.DataFrame())
    inv = tables.get("inventory_diagnostics", pd.DataFrame())
    quotes = tables.get("quote_diagnostics", pd.DataFrame())
    market = tables.get("market_quality", pd.DataFrame())
    markout = tables.get("markout_diagnostics", pd.DataFrame())
    risk = tables.get("risk_diagnostics", pd.DataFrame())
    fills = tables.get("fill_diagnostics", pd.DataFrame())

    pnl_ts = _read_csv(Path(diagnostics["artifacts"].get("pnl_timeseries.csv", "")), warnings=diagnostics["warnings"], optional=True)
    if not pnl_ts.empty and "total_pnl" in pnl_ts:
        _coerce_numeric(pnl_ts)
        save_plot("pnl_curve.png", lambda plt: plt.plot(pnl_ts["total_pnl"]))
        save_plot("drawdown.png", lambda plt: plt.plot(pnl_ts["total_pnl"] - pnl_ts["total_pnl"].cummax()))
    if not inv.empty and "inventory" in inv:
        save_plot("inventory_timeseries.png", lambda plt: plt.plot(inv["inventory"]))
    if not quotes.empty and "quote_spread_bps" in quotes:
        save_plot("quoted_spread_distribution.png", lambda plt: plt.hist(quotes["quote_spread_bps"].dropna(), bins=30))
    if not market.empty and "spread_bps" in market:
        save_plot("book_spread_distribution.png", lambda plt: plt.hist(market["spread_bps"].dropna(), bins=30))
        save_plot("orderbook_quality.png", lambda plt: market[["is_crossed_book", "is_missing_top_of_book", "is_wide_spread"]].sum().plot(kind="bar", ax=plt.gca()))
    if not market.empty and "imbalance_1" in market:
        save_plot("imbalance_distribution.png", lambda plt: plt.hist(market["imbalance_1"].dropna(), bins=30))
    if not markout.empty:
        cols = [c for c in markout.columns if c.startswith("markout_bps_h")]
        if cols:
            save_plot("fill_markout_distribution.png", lambda plt: plt.hist(markout[cols[0]].dropna(), bins=30))
    save_plot("event_funnel.png", lambda plt: plt.bar(["input", "quoted", "fills"], [diagnostics["run"].get("input_events") or 0, diagnostics["run"].get("quoted_events") or 0, diagnostics["run"].get("number_of_fills") or 0]))
    if not fills.empty and "side" in fills:
        save_plot("fill_side_counts.png", lambda plt: fills["side"].value_counts().plot(kind="bar", ax=plt.gca()))
    if not risk.empty:
        save_plot("risk_reason_counts.png", lambda plt: risk.set_index("reason")["count"].plot(kind="bar", ax=plt.gca()))
    save_plot("pnl_attribution.png", lambda plt: plt.bar(["net", "fees"], [diagnostics["pnl"].get("net_pnl") or 0, diagnostics["pnl"].get("fees") or 0]))


# Utility helpers
def _read_json(path: Path, *, warnings: list[str]) -> dict[str, Any]:
    if not path.exists():
        warnings.append(f"Missing summary file: {path.name}")
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception as exc:
        warnings.append(f"Could not read {path}: {exc}")
        return {}


def _read_csv(path: Path, *, warnings: list[str], optional: bool) -> pd.DataFrame:
    if not path or not path.exists():
        if not optional:
            warnings.append(f"Missing CSV file: {path}")
        return pd.DataFrame()
    try:
        return pd.read_csv(path)
    except pd.errors.EmptyDataError:
        return pd.DataFrame()
    except Exception as exc:
        warnings.append(f"Could not read {path}: {exc}")
        return pd.DataFrame()


def _coerce_numeric(frame: pd.DataFrame) -> None:
    for col in frame.columns:
        if col not in {"timestamp", "symbol", "side", "order_id", "event_type", "quote_reason", "risk_reason"}:
            converted = pd.to_numeric(frame[col], errors="coerce")
            if converted.notna().any() or frame[col].isna().all():
                frame[col] = converted


def _to_datetime(series: Any) -> pd.Series:
    if series is None:
        return pd.Series(dtype="datetime64[ns, UTC]")
    return pd.to_datetime(series, errors="coerce", utc=True)


def _to_bool(series: pd.Series) -> pd.Series:
    if series.dtype == bool:
        return series
    return series.astype(str).str.lower().isin({"true", "1", "yes"})


def _mean(series: Any) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if s.empty else float(s.mean())


def _median(series: Any) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if s.empty else float(s.median())


def _quantile(series: Any, q: float) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if s.empty else float(s.quantile(q))


def _min(series: Any) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if s.empty else float(s.min())


def _max(series: Any) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if s.empty else float(s.max())


def _std(series: Any) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if len(s) < 2 else float(s.std(ddof=1))


def _sum(series: Any) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if s.empty else float(s.sum())


def _last(series: Any) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if s.empty else float(s.iloc[-1])


def _delta(series: Any) -> float | None:
    if series is None:
        return None
    s = pd.Series(series).dropna()
    return None if len(s) < 2 else float(s.iloc[-1] - s.iloc[0])


def _safe_div(num: Any, den: Any) -> float | None:
    try:
        if num is None or den is None or float(den) == 0.0:
            return None
        return float(num) / float(den)
    except Exception:
        return None


def _counts(series: Any) -> dict[str, int]:
    if series is None:
        return {}
    s = pd.Series(series).dropna().astype(str)
    return {k: int(v) for k, v in s.value_counts().to_dict().items()}


def _contains_count(series: pd.Series, pattern: str) -> int:
    return int(series.astype(str).str.contains(pattern, case=False, regex=False).sum())


def _coalesce(*values: Any) -> Any:
    for value in values:
        if value is not None and not (isinstance(value, float) and np.isnan(value)):
            return value
    return None


def _run_bounds(frames: Sequence[pd.DataFrame]) -> tuple[pd.Timestamp | None, pd.Timestamp | None]:
    timestamps: list[pd.Timestamp] = []
    for frame in frames:
        if not frame.empty and "timestamp" in frame:
            ts = _to_datetime(frame["timestamp"]).dropna()
            if not ts.empty:
                timestamps.extend([ts.min(), ts.max()])
    if not timestamps:
        return None, None
    return min(timestamps), max(timestamps)


def _jsonable(value: Any) -> Any:
    if isinstance(value, dict):
        return {str(k): _jsonable(v) for k, v in value.items() if k != "_tables"}
    if isinstance(value, list):
        return [_jsonable(v) for v in value]
    if isinstance(value, (np.integer, np.floating)):
        if pd.isna(value):
            return None
        return value.item()
    if isinstance(value, pd.Timestamp):
        return value.isoformat()
    if isinstance(value, float) and np.isnan(value):
        return None
    return value


def _comparison_markdown(frame: pd.DataFrame) -> str:
    lines = ["# Market Making Run Comparison", ""]
    if frame.empty:
        lines.append("No runs discovered.")
    else:
        cols = list(frame.columns)
        lines.append("| " + " | ".join(cols) + " |")
        lines.append("| " + " | ".join(["---"] * len(cols)) + " |")
        for _, row in frame.iterrows():
            lines.append("| " + " | ".join(str(row[col]) for col in cols) + " |")
    lines.append("")
    return "\n".join(lines)


def _write_comparison_pptx(frame: pd.DataFrame, output_path: Path) -> None:
    try:
        from pptx import Presentation
    except Exception:
        return
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "Market Making Run Comparison"
    title.placeholders[1].text = f"Runs analyzed: {len(frame)}"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, row in frame.head(8).iterrows():
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = f"{row.get('run_dir')}: PnL={row.get('total_pnl')}, fills={row.get('fills')}, fill_ratio={row.get('fill_ratio')}"
    prs.save(output_path)


__all__ = [
    "build_market_making_diagnostics",
    "discover_market_making_runs",
    "write_market_making_comparison",
    "write_market_making_diagnostics",
]
